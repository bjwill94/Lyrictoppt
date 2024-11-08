import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from io import BytesIO

# Path to the sample PPT file
sample_ppt_path = "sample.pptx"

# Function to add song lyrics as new slides with formatting
def add_lyrics_to_ppt(ppt, lyrics):
    # Split the lyrics into paragraphs by detecting blank lines
    paragraphs = [p.strip() for p in lyrics.strip().split('\n\n') if p.strip()]

    # Initialize a counter for slide numbering
    slide_number = 1

    # Loop through each paragraph and add it as a new slide
    for paragraph in paragraphs:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # 1 for a Title and Content layout
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]

        # Set the dynamic title of the slide
        title.text = f"Song Lyrics {slide_number}"
        title.text_frame.paragraphs[0].font.size = Pt(18)  # Set title font size to 18

        # Set the body (text box) content with the paragraph
        text_frame = body.text_frame
        text_frame.clear()  # Clear any default text

        # Add the paragraph text without bullet points
        p = text_frame.add_paragraph()
        p.text = paragraph

        # Ensure no bullet points
        p.level = 0  # Level 0 is no bullet point
        p.bullet = None  # Explicitly remove bullets

        # Center align the text
        p.alignment = PP_ALIGN.CENTER  # Center alignment horizontally
        # text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center alignment vertically

        # Set a larger font size for the lyrics text
        p.font.size = Pt(24)  # Font size for lyrics set to 24

        # Set the line spacing to 1.5 directly on the paragraph
        p.line_spacing = 1.5  # Setting line spacing directly on paragraph

        # Increment the slide number
        slide_number += 1

    return ppt

# Function to convert the PPT to a downloadable file
def convert_ppt_to_bytes(ppt):
    ppt_bytes_io = BytesIO()
    ppt.save(ppt_bytes_io)
    ppt_bytes_io.seek(0)
    return ppt_bytes_io

# Streamlit interface
st.title('Lyrics to PPT Converter')

# Upload PPT file
uploaded_ppt = st.file_uploader("Upload PowerPoint (.pptx) file", type="pptx")

# Check if a file was uploaded, otherwise use the sample PPT
if uploaded_ppt:
    ppt = Presentation(uploaded_ppt)
    st.write("Using uploaded PowerPoint file.")
else:
    # Load the sample PPT if no file is uploaded
    ppt = Presentation(sample_ppt_path)
    st.write("No PowerPoint uploaded. Using the sample PowerPoint file instead.")

# Input song name
song_name = st.text_input("Enter Song Name (used as PPT file name)")

# Input song lyrics
lyrics = st.text_area("Enter Song Lyrics", height=300)

if song_name and lyrics:
    # Process to add lyrics into the PPT
    if st.button("Generate PPT with Lyrics"):
        updated_ppt = add_lyrics_to_ppt(ppt, lyrics)
        ppt_bytes_io = convert_ppt_to_bytes(updated_ppt)

        # Provide a download link for the new PPT, using song_name as file name
        st.download_button(
            label="Download Updated PPT",
            data=ppt_bytes_io,
            file_name=f"{song_name}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

