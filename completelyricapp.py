import streamlit as st
import webbrowser
import urllib.parse
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from io import BytesIO

# Function to generate the search query for lyrics
def generate_search_query(lyrics, language):
    # Construct the search query as "lyrics + language + christian devotional song"
    search_query = f"lyrics {lyrics} {language} christian devotional"
    # URL encode the search query
    encoded_query = urllib.parse.quote_plus(search_query)
    # Return the complete search URL (Google search)
    return f"https://www.google.com/search?q={encoded_query}"

# Function to add song lyrics to PPT
def add_lyrics_to_ppt(ppt, lyrics):
    # Split the lyrics into paragraphs by detecting blank lines
    paragraphs = [p.strip() for p in lyrics.strip().split('\n\n') if p.strip()]

    # Initialize a counter for slide numbering
    slide_number = 1

    # Loop through each paragraph and add it as a new slide
    for paragraph in paragraphs:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # 1 for a Title and Content layout
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]

        # Set the dynamic title of the slide
        title.text = f"Song Lyrics {slide_number}"
        title.text_frame.paragraphs[0].font.size = Pt(18)  # Set title font size to 18

        # Set the body (text box) content with the paragraph
        text_frame = body.text_frame
        text_frame.clear()  # Clear any default text

        # Add the paragraph text without bullet points
        p = text_frame.add_paragraph()
        p.text = paragraph

        # Ensure no bullet points
        p.level = 0  # Level 0 is no bullet point
        p.bullet = None  # Explicitly remove bullets

        # Center align the text
        p.alignment = PP_ALIGN.CENTER  # Center alignment horizontally
        # text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center alignment vertically

        # Set a larger font size for the lyrics text
        p.font.size = Pt(24)  # Font size for lyrics set to 24

        # Set the line spacing to 1.5 directly on the paragraph
        p.line_spacing = 1.5  # Setting line spacing directly on paragraph

        # Increment the slide number
        slide_number += 1

    return ppt

# Function to convert the PPT to a downloadable file
def convert_ppt_to_bytes(ppt):
    ppt_bytes_io = BytesIO()
    ppt.save(ppt_bytes_io)
    ppt_bytes_io.seek(0)
    return ppt_bytes_io

# Streamlit interface
def lyrics_searcher_page():
    # Input for first line of lyrics
    lyrics_input = st.text_input("Enter the first line of the lyrics:")

    # Input for language
    language_options = ["English", "Malayalam", "Tamil", "Hindi"]
    language_input = st.selectbox("Select the language:", language_options)

    # Button to perform search
    if st.button("Search"):
        if lyrics_input:
            # Generate the search query based on the inputs
            search_url = generate_search_query(lyrics_input, language_input)

            # Display the generated search URL
            st.write(f"Searching for: {search_url}")

            # Open the first search result link in a new tab
            webbrowser.open_new_tab(search_url)
        else:
            st.error("Please enter the first line of lyrics.")

    # Button to navigate to the PPT Generator page
    if st.button("Go to PPT Generator"):
        st.session_state.page = "PPT Generator"

def ppt_generator_page():
    st.title("Lyrics to PPT Converter")

    # Upload PPT file
    uploaded_ppt = st.file_uploader("Upload PowerPoint (.pptx) file", type="pptx")
    if uploaded_ppt:
        ppt = Presentation(uploaded_ppt)

        # Input song name
        song_name = st.text_input("Enter Song Name (used as PPT file name)")

        # Input song lyrics
        lyrics = st.text_area("Enter Song Lyrics", height=300)

        if song_name and lyrics:
            # Process to add lyrics into the PPT
            if st.button("Generate PPT with Lyrics"):
                updated_ppt = add_lyrics_to_ppt(ppt, lyrics)
                ppt_bytes_io = convert_ppt_to_bytes(updated_ppt)

                # Provide a download link for the new PPT, using song_name as file name
                st.download_button(
                    label="Download Updated PPT",
                    data=ppt_bytes_io,
                    file_name=f"{song_name}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

# Main app logic
if "page" not in st.session_state:
    st.session_state.page = "Lyrics Searcher"

# Page navigation logic
if st.session_state.page == "Lyrics Searcher":
    lyrics_searcher_page()
elif st.session_state.page == "PPT Generator":
    ppt_generator_page()
